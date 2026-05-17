from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ──────────────────────────────────────────────────
NAVY   = RGBColor(26,  46,  74)
TEAL   = RGBColor(0,  150, 160)
AMBER  = RGBColor(244,168,  32)
GREEN  = RGBColor(0,  140,  80)
LIGHT  = RGBColor(242,246, 250)
WHITE  = RGBColor(255,255, 255)
DARK   = RGBColor(30,  30,  30)
MID    = RGBColor(100,110, 125)
RED    = RGBColor(200,  30,  30)

IMGS = r"C:\Users\romai\Desktop\saillowtech\Report_Sailowtech\images"
def img(n): return os.path.join(IMGS, n)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height
BL = prs.slide_layouts[6]   # blank


# ── Helpers ──────────────────────────────────────────────────
def rect(slide, l, t, w, h, color, line=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s

def tx(slide, l, t, w, h, text, size=16, bold=False, color=DARK,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text          = text
    run.font.size     = Pt(size)
    run.font.bold     = bold
    run.font.italic   = italic
    run.font.color.rgb = color
    return box

def picture(slide, fname, l, t, w, h):
    try:
        slide.shapes.add_picture(img(fname), l, t, w, h)
    except Exception as e:
        print(f"  [img skip] {fname}: {e}")

def title_bar(slide, text, accent=TEAL):
    rect(slide, 0, 0, W, Cm(2.2), NAVY)
    rect(slide, 0, Cm(2.2), W, Cm(0.15), accent)
    rect(slide, 0, Cm(2.35), W, H - Cm(2.35), LIGHT)
    tx(slide, Cm(0.9), Cm(0.3), Cm(30), Cm(1.75),
       text, size=24, bold=True, color=WHITE)

def bullet_list(slide, l, t, w, h, items, size=17, color=DARK, indent="• "):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(7)
        run = p.add_run()
        run.text = indent + item
        run.font.size = Pt(size)
        run.font.color.rgb = color

def divider(title, subtitle, img_file=None):
    s = prs.slides.add_slide(BL)
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, Cm(0.5), H, TEAL)
    rect(s, 0, H - Cm(0.3), W, Cm(0.3), AMBER)
    if img_file:
        picture(s, img_file, Inches(7.8), Cm(1.5), Inches(5.0), Inches(5.2))
        rect(s, Inches(7.8), Cm(1.5), Inches(5.0), Inches(5.2),
             RGBColor(10, 28, 52))  # tint overlay
    tx(s, Cm(1.5), Cm(3.2), Cm(28), Cm(2.5),
       title, size=40, bold=True, color=WHITE)
    tx(s, Cm(1.5), Cm(5.9), Cm(22), Cm(1.2),
       subtitle, size=20, italic=True, color=TEAL)
    return s


# ════════════════════════════════════════════════════════════
# SLIDE 1 — TITRE
# ════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BL)
rect(s1, 0, 0, W, H, NAVY)

# Photo voilier à droite
picture(s1, "voilier_sunset.png", Inches(7.5), 0, Inches(5.83), Inches(7.5))
rect(s1, Inches(7.5), 0, Inches(5.83), H, RGBColor(10, 28, 52))  # overlay

# Barre teal verticale séparatrice
rect(s1, Inches(7.35), 0, Cm(0.45), H, TEAL)
# Barre amber haute gauche
rect(s1, 0, 0, Inches(7.35), Cm(0.35), AMBER)

# Logo EPFL
picture(s1, "epfl.png", Cm(0.9), Cm(0.7), Cm(3.2), Cm(1.1))

# Titre principal
tx(s1, Cm(0.9), Cm(3.8), Cm(26), Cm(3.0),
   "Station météo\net Métadonnées", size=36, bold=True, color=WHITE)

# Ligne teal
rect(s1, Cm(0.9), Cm(7.2), Cm(10), Cm(0.12), TEAL)

# Sous-titre
tx(s1, Cm(0.9), Cm(7.5), Cm(24), Cm(1.0),
   "Projet Sailowtech", size=22, color=TEAL)

# Infos bas
tx(s1, Cm(0.9), Cm(9.4), Cm(24), Cm(0.8),
   "Guilhem Destriau  ·  Romain Corbel", size=16,
   color=RGBColor(180, 200, 220))
tx(s1, Cm(0.9), Cm(10.2), Cm(24), Cm(0.7),
   "Superviseur : Eric Boillat  —  LMTM, EPFL", size=14,
   color=RGBColor(150, 175, 200))
tx(s1, Cm(0.9), Cm(11.0), Cm(24), Cm(0.6),
   "ME-401  —  Lausanne, Mai 2026", size=13,
   color=RGBColor(130, 155, 180))


# ════════════════════════════════════════════════════════════
# SLIDE 2 — SOMMAIRE
# ════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BL)
title_bar(s2, "Sommaire", AMBER)

items_sommaire = [
    ("01", "Contexte : le projet Sailowtech"),
    ("02", "Objectifs du projet"),
    ("03", "Design mécanique"),
    ("04", "Conception électronique"),
    ("05", "Vision par ordinateur"),
    ("06", "Résultats et perspectives"),
]
for i, (num, label) in enumerate(items_sommaire):
    y = Cm(2.9) + i * Cm(1.45)
    rect(s2, Cm(0.9), y, Cm(1.6), Cm(1.2), NAVY)
    tx(s2, Cm(0.9), y + Cm(0.15), Cm(1.6), Cm(0.95),
       num, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s2, Cm(2.9), y + Cm(0.2), Cm(25), Cm(0.85),
       label, size=20, color=DARK)


# ════════════════════════════════════════════════════════════
# SLIDE 3 — CONTEXTE SAILOWTECH
# ════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BL)
title_bar(s3, "Contexte : le projet Sailowtech")

bullet_list(s3, Cm(0.9), Cm(2.8), Cm(18), Cm(10.5), [
    "Projet MAKE de l'EPFL — association fondée sur l'open-source et le low-tech",
    "Expéditions scientifiques à la voile : lacs, mers et océans",
    "Objectif : transformer des voiliers en plateformes de collecte de données environnementales",
    "Approche frugale et participative — instruments accessibles, données partagées librement",
    "Grandeurs visées : qualité de l'air, paramètres de la colonne d'eau, rayonnement solaire",
], size=18)

picture(s3, "voilier_sunset.png", Cm(20.5), Cm(2.7), Cm(12.5), Cm(9.0))
rect(s3, Cm(20.5), Cm(11.7), Cm(12.5), Cm(0.8), NAVY)
tx(s3, Cm(20.5), Cm(11.75), Cm(12.5), Cm(0.7),
   "Expédition Sailowtech", size=12, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 4 — OBJECTIFS
# ════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BL)
title_bar(s4, "Objectifs du projet", TEAL)

objs = [
    (TEAL,  "Données atmosphériques",         "PM₁.₀ / PM₂.₅ / PM₁₀, irradiance, T°, pression, humidité"),
    (AMBER, "Réseau du bateau (NMEA 2000)",   "GPS, vitesse, cap, profondeur depuis l'ordinateur de bord (ex. Garmin)"),
    (GREEN, "Sonde CTD",                      "Température, conductivité, profondeur — RS-485 câblé ou Zigbee"),
    (NAVY,  "Résistance marine",              "Projections d'eau, corrosion saline, vibrations"),
    (TEAL,  "Autonomie énergétique",          "Batterie 12 V avec régulation DC-DC embarquée"),
    (RED,   "Vision par ordinateur (bonus)",  "Détection automatique de la couverture nuageuse par CNN"),
]
for i, (col, title, desc) in enumerate(objs):
    y = Cm(2.7) + i * Cm(1.5)
    rect(s4, Cm(0.5), y + Cm(0.25), Cm(0.35), Cm(0.85), col)
    tx(s4, Cm(1.2), y,          Cm(18), Cm(0.72), title, size=16, bold=True, color=col)
    tx(s4, Cm(1.2), y + Cm(0.68), Cm(18), Cm(0.72), desc, size=13, color=MID)

picture(s4, "station_assemblée.png", Cm(20.5), Cm(2.5), Cm(12.5), Cm(11.5))


# ════════════════════════════════════════════════════════════
# SLIDE 5 — SECTION : DESIGN MÉCANIQUE
# ════════════════════════════════════════════════════════════
divider("Design mécanique",
        "Modularité  ·  Exploration des designs  ·  Abri de Stevenson",
        "mecanique_boitier_photo.png")


# ════════════════════════════════════════════════════════════
# SLIDE 6 — MODULARITÉ
# ════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BL)
title_bar(s6, "Design mécanique — Principe de modularité", TEAL)

# Encadré principe clé
rect(s6, Cm(0.9), Cm(2.65), Cm(18.5), Cm(1.3), NAVY)
tx(s6, Cm(1.2), Cm(2.75), Cm(18.0), Cm(1.1),
   "Principe fondateur : aucune pièce n'est fixée définitivement",
   size=19, bold=True, color=AMBER)

bullet_list(s6, Cm(0.9), Cm(4.2), Cm(18.5), Cm(9.0), [
    "Vis, entretoises et clips permettent la modification sans outil spécial",
    "Hauteur augmentable pour accueillir capteurs supplémentaires ou batteries",
    "Chaque étage PCB est démontable indépendamment — maintenance facilitée",
    "PCB circulaires : optimisation de l'espace dans le boîtier cylindrique",
    "Connecteurs inter-étages positionnés pour minimiser la longueur des câbles",
    "Prototypes intérieurs validés par découpe laser sur MDF avant commande des PCB",
], size=17)

picture(s6, "mecanique_interieur_cao.png",  Cm(20.5), Cm(2.6), Cm(12.5), Cm(5.5))
picture(s6, "mecanique_interieur_photo.png", Cm(20.5), Cm(8.4), Cm(12.5), Cm(5.3))


# ════════════════════════════════════════════════════════════
# SLIDE 7 — EXPLORATION DES DESIGNS
# ════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BL)
title_bar(s7, "Boîtier extérieur — Exploration des designs", AMBER)

steps = [
    (RED,   "1",
     "Boîtier cylindrique fermé",
     "Design initial hérité des projets précédents.\n"
     "Sans ventilation active : renouvellement d'air insuffisant."),
    (RGBColor(180, 100, 0), "2",
     "Membranes Gore-Tex",
     "L'air passe sans que l'eau ne pénètre.\n"
     "Mais la résistance à l'écoulement est trop élevée :\n"
     "ventilation passive insuffisante."),
    (GREEN, "3",
     "Abri de Stevenson cylindrique",
     "Différentes longueurs et inclinaisons testées par\n"
     "impression 3D. Design final : meilleur compromis\n"
     "entre protection et aération."),
]
for i, (col, num, title, desc) in enumerate(steps):
    x = Cm(0.6) + i * Cm(11.2)
    rect(s7, x, Cm(2.7), Cm(10.7), Cm(1.2), col)
    tx(s7, x + Cm(0.3), Cm(2.8), Cm(1.3), Cm(1.0),
       num, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s7, x + Cm(1.8), Cm(2.85), Cm(8.6), Cm(1.0),
       title, size=16, bold=True, color=WHITE)
    tx(s7, x + Cm(0.3), Cm(4.1), Cm(10.2), Cm(3.2),
       desc, size=14, color=DARK, wrap=True)

picture(s7, "mecanique_boitier_cao.png",       Cm(0.6),  Cm(7.2), Cm(10.7), Cm(6.4))
picture(s7, "mecanique_stevenson_anneau.png",  Cm(11.8), Cm(7.2), Cm(10.7), Cm(6.4))
picture(s7, "mecanique_boitier_photo.png",     Cm(23.0), Cm(7.2), Cm(10.7), Cm(6.4))


# ════════════════════════════════════════════════════════════
# SLIDE 8 — STATION ASSEMBLÉE
# ════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BL)
title_bar(s8, "Station assemblée et validée", GREEN)

facts = [
    (TEAL,  "Entièrement blanche",
     "Limite l'absorption solaire et maintient la T° interne proche de l'ambiante"),
    (AMBER, "Plexiglas réduit au strict minimum",
     "Seule la fenêtre caméra est transparente — effet de serre maîtrisé"),
    (GREEN, "Tests d'étanchéité et de résistance réalisés",
     "Aspersions d'eau et chocs simulés : assemblage et joints validés"),
    (NAVY,  "Validé en conditions réelles",
     "Design polyvalent, conçu pour toutes les conditions non extrêmes"),
]
for i, (col, title, desc) in enumerate(facts):
    y = Cm(2.8) + i * Cm(1.9)
    rect(s8, Cm(0.5), y + Cm(0.3), Cm(0.4), Cm(0.9), col)
    tx(s8, Cm(1.2), y,            Cm(19), Cm(0.78), title, size=18, bold=True, color=col)
    tx(s8, Cm(1.2), y + Cm(0.75), Cm(19), Cm(0.85), desc,  size=14, color=MID)

picture(s8, "station_assemblée.png", Cm(21.0), Cm(2.5), Cm(12.3), Cm(11.5))


# ════════════════════════════════════════════════════════════
# SLIDE 9 — SECTION : CONCEPTION ÉLECTRONIQUE
# ════════════════════════════════════════════════════════════
divider("Conception électronique",
        "Architecture  ·  PCB étage 1  ·  PCB étage 2",
        "pcb1_3d.png")


# ════════════════════════════════════════════════════════════
# SLIDE 10 — ARCHITECTURE DE COMMUNICATION
# ════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BL)
title_bar(s10, "Architecture de communication")

protos = [
    (TEAL,  "NMEA 2000 (CAN Bus)",   "Réseau du bateau — GPS, vitesse, cap depuis l'ordinateur de bord (ex. Garmin)"),
    (AMBER, "RS-485 / RS-232",       "Sonde CTD câblée — communication continue et en temps réel"),
    (GREEN, "Zigbee (XIAO ESP32C6)", "Sonde CTD sans câble — échange de paramètres avant immersion"),
    (NAVY,  "Wi-Fi (ESP32)",         "Transmission vers l'ordinateur portable à bord — visualisation temps réel"),
    (RGBColor(120,60,180),
              "SPI / I²C / UART",    "Communication interne entre microcontrôleur et capteurs / modules"),
]
for i, (col, proto, desc) in enumerate(protos):
    y = Cm(2.7) + i * Cm(1.6)
    rect(s10, Cm(0.5), y, Cm(9.0), Cm(1.35), col)
    tx(s10, Cm(0.8), y + Cm(0.2), Cm(8.5), Cm(1.0),
       proto, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s10, Cm(10.0), y + Cm(0.2), Cm(11.5), Cm(1.1), desc, size=14, color=DARK)

picture(s10, "protocole_comm.png", Cm(22.5), Cm(2.5), Cm(11.0), Cm(11.0))


# ════════════════════════════════════════════════════════════
# SLIDE 11 — PCB ÉTAGE 1
# ════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BL)
title_bar(s11, "PCB étage 1 — Contrôle et connectivité  (situé en bas de la station)", TEAL)

comps11 = [
    ("ESP32 DevKit M V1",  "Microcontrôleur principal — Wi-Fi, BT, 240 MHz, 2 cœurs"),
    ("XIAO ESP32C6",       "Module Zigbee — échange avec la sonde CTD sans câble"),
    ("MCP2515",            "Transceiver CAN — interface NMEA 2000"),
    ("Adafruit GPS",       "Horodatage et position géographique de chaque mesure"),
    ("Adafruit micro SD",  "Enregistrement local des données en CSV"),
    ("LM2596",             "Régulateur buck 12 V → 5 V"),
    ("Ports CTD & NMEA",   "Installés sur le PCB — perçages à réaliser lors des tests terrain"),
    ("Ports supp.",        "Réservés pour l'ajout de capteurs futurs (modularité)"),
]
for i, (comp, desc) in enumerate(comps11):
    y = Cm(2.7) + i * Cm(1.3)
    rect(s11, Cm(0.5), y + Cm(0.3), Cm(0.25), Cm(0.6), TEAL)
    tx(s11, Cm(1.0), y,             Cm(10.5), Cm(0.7), comp, size=14, bold=True, color=NAVY)
    tx(s11, Cm(1.0), y + Cm(0.65), Cm(17.5), Cm(0.62), desc, size=12, color=MID)

picture(s11, "pcb1_routage.png", Cm(20.0), Cm(2.7), Cm(6.5), Cm(6.0))
picture(s11, "pcb1_3d.png",      Cm(27.0), Cm(2.7), Cm(6.3), Cm(6.0))

rect(s11, Cm(0.5), H - Cm(1.4), Cm(19.0), Cm(1.15), RGBColor(228, 236, 248))
tx(s11, Cm(0.8), H - Cm(1.35), Cm(18.5), Cm(1.05),
   "Les connexions CTD et NMEA sont installées sur le PCB. Les perçages dans le boîtier mécanique "
   "seront réalisés lors de la validation avec un ordinateur de bord.",
   size=11, italic=True, color=MID)


# ════════════════════════════════════════════════════════════
# SLIDE 12 — PCB ÉTAGE 2
# ════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(BL)
title_bar(s12, "PCB étage 2 — Capteurs  (exposé à l'air ambiant)", AMBER)

comps12 = [
    ("SPS30 (Sensirion)",    "PM₁.₀ / PM₂.₅ / PM₄ / PM₁₀ — diffraction laser"),
    ("BME280",               "Température ±1 °C, Pression ±1 hPa, Humidité ±3 %"),
    ("BH1750",               "Éclairement lumineux en lux — indicateur d'irradiance solaire"),
    ("MPU-6050",             "IMU 6 axes — correction tangage / roulis pour l'irradiance"),
    ("Joy-it ESP32 camera",  "Images zénithales — alimentation de l'algorithme de vision"),
    ("Ports supplémentaires","Réservés pour l'ajout de capteurs futurs (modularité)"),
]
for i, (comp, desc) in enumerate(comps12):
    y = Cm(2.7) + i * Cm(1.5)
    rect(s12, Cm(0.5), y + Cm(0.3), Cm(0.25), Cm(0.7), AMBER)
    tx(s12, Cm(1.0), y,             Cm(18), Cm(0.78), comp, size=15, bold=True, color=NAVY)
    tx(s12, Cm(1.0), y + Cm(0.75), Cm(18), Cm(0.72), desc, size=13, color=MID)

picture(s12, "pcb2_routage.png", Cm(20.0), Cm(2.7), Cm(6.5), Cm(6.0))
picture(s12, "pcb2_3d.png",      Cm(27.0), Cm(2.7), Cm(6.3), Cm(6.0))

rect(s12, Cm(0.5), H - Cm(1.4), Cm(19.0), Cm(1.15), RGBColor(228, 236, 248))
tx(s12, Cm(0.8), H - Cm(1.35), Cm(18.5), Cm(1.05),
   "Cet étage est positionné dans la zone ventilée du boîtier (abri de Stevenson) "
   "pour éviter tout biais thermique sur les mesures atmosphériques.",
   size=11, italic=True, color=MID)


# ════════════════════════════════════════════════════════════
# SLIDE 13 — SECTION : VISION PAR ORDINATEUR
# ════════════════════════════════════════════════════════════
divider("Vision par ordinateur",
        "Segmentation sémantique  ·  U-Net  ·  Couverture nuageuse",
        "vision_ciel.png")


# ════════════════════════════════════════════════════════════
# SLIDE 14 — SEGMENTATION NUAGEUSE
# ════════════════════════════════════════════════════════════
s14 = prs.slides.add_slide(BL)
title_bar(s14, "Détection automatique de la couverture nuageuse")

arch = [
    ("Tâche",      "Segmentation sémantique binaire pixel par pixel — nuage (1) / ciel (0)"),
    ("Modèle",     "U-Net  +  encodeur ResNet-34 pré-entraîné sur ImageNet"),
    ("Dataset",    "SWIMSEG — 861 images d'entraînement · 101 validation · 51 test"),
    ("Optimiseur", "AdamW  lr=1e-4 · CosineAnnealingLR · 30 époques · batch size 8"),
    ("Fonction de coût", "L = 0.5 × BCE  +  0.5 × Dice Loss"),
]
for i, (label, val) in enumerate(arch):
    y = Cm(2.7) + i * Cm(1.6)
    rect(s14, Cm(0.5), y, Cm(5.2), Cm(1.35), NAVY)
    tx(s14, Cm(0.65), y + Cm(0.22), Cm(4.9), Cm(0.95),
       label, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s14, Cm(6.2), y + Cm(0.22), Cm(13.5), Cm(1.05), val, size=14, color=DARK)

picture(s14, "vision_ciel.png",   Cm(20.8), Cm(2.6), Cm(12.5), Cm(5.8))
picture(s14, "vision_masque.png", Cm(20.8), Cm(8.7), Cm(12.5), Cm(5.8))
tx(s14, Cm(20.8), Cm(8.4), Cm(12.5), Cm(0.5),
   "Image brute du ciel                      Masque prédit (blanc = nuage)",
   size=10, color=MID, italic=True)


# ════════════════════════════════════════════════════════════
# SLIDE 15 — RÉSULTATS
# ════════════════════════════════════════════════════════════
s15 = prs.slides.add_slide(BL)
title_bar(s15, "Résultats", GREEN)

# Colonne gauche
rect(s15, Cm(0.5), Cm(2.6), Cm(16.5), Cm(1.0), NAVY)
tx(s15, Cm(0.8), Cm(2.65), Cm(16.0), Cm(0.88),
   "Électronique et firmware", size=17, bold=True, color=WHITE)
for i, r in enumerate([
    "Tous les capteurs testés individuellement et en configuration complète",
    "Communication inter-PCB stable · SD fonctionnel · Wi-Fi opérationnel",
    "Firmware complet : acquisition, horodatage, stockage, transmission",
]):
    tx(s15, Cm(1.0), Cm(3.75) + i * Cm(0.85), Cm(15.5), Cm(0.78),
       "✓  " + r, size=14, color=DARK)

rect(s15, Cm(0.5), Cm(6.5), Cm(16.5), Cm(1.0), GREEN)
tx(s15, Cm(0.8), Cm(6.55), Cm(16.0), Cm(0.88),
   "Validation mécanique", size=17, bold=True, color=WHITE)
for i, r in enumerate([
    "Étanchéité validée — aspersions d'eau et joints testés",
    "Résistance mécanique confirmée (chocs simulés)",
    "Tests en conditions réelles effectués — design approuvé",
]):
    tx(s15, Cm(1.0), Cm(7.65) + i * Cm(0.85), Cm(15.5), Cm(0.78),
       "✓  " + r, size=14, color=DARK)

# Colonne droite — tableau vision
rect(s15, Cm(18.5), Cm(2.6), Cm(15.0), Cm(1.0), TEAL)
tx(s15, Cm(18.8), Cm(2.65), Cm(14.5), Cm(0.88),
   "Vision par ordinateur — SWIMSEG (51 images test)", size=16, bold=True, color=WHITE)

metrics = [("Métrique", "Valeur", NAVY, True),
           ("IoU moyen",  "0.888",  LIGHT, False),
           ("Dice moyen", "0.938",  WHITE, False),
           ("F1 moyen",   "0.935",  LIGHT, False)]
for i, (met, val, bg, header) in enumerate(metrics):
    y = Cm(3.7) + i * Cm(0.95)
    rect(s15, Cm(18.5), y, Cm(9.5), Cm(0.88), bg if not header else NAVY)
    rect(s15, Cm(28.0), y, Cm(5.5), Cm(0.88), bg if not header else NAVY)
    c = WHITE if header else DARK
    tx(s15, Cm(18.8), y + Cm(0.07), Cm(9.0), Cm(0.75),
       met, size=14, bold=header, color=c)
    tx(s15, Cm(28.1), y + Cm(0.07), Cm(5.2), Cm(0.75),
       val, size=15 if not header else 14, bold=True, color=TEAL if not header else WHITE,
       align=PP_ALIGN.CENTER)

rect(s15, Cm(18.5), Cm(7.6), Cm(15.0), Cm(1.3), RGBColor(0, 130, 140))
tx(s15, Cm(18.8), Cm(7.7), Cm(14.5), Cm(1.1),
   "F1 = 0.935 en moyenne\nMinimum 0.815 sur images difficiles",
   size=15, bold=True, color=WHITE)


# ════════════════════════════════════════════════════════════
# SLIDE 16 — CONCLUSION & PERSPECTIVES
# ════════════════════════════════════════════════════════════
s16 = prs.slides.add_slide(BL)
rect(s16, 0, 0, W, H, NAVY)
rect(s16, 0, 0, Cm(0.5), H, AMBER)
rect(s16, 0, H - Cm(0.3), W, Cm(0.3), TEAL)

tx(s16, Cm(1.5), Cm(0.7), Cm(30), Cm(1.6),
   "Conclusion et perspectives", size=32, bold=True, color=WHITE)
rect(s16, Cm(1.5), Cm(2.4), Cm(13), Cm(0.12), AMBER)

# Réalisé
rect(s16, Cm(1.5), Cm(2.65), Cm(14.5), Cm(1.0), TEAL)
tx(s16, Cm(1.8), Cm(2.7), Cm(14.0), Cm(0.88),
   "Ce qui a été réalisé", size=18, bold=True, color=WHITE)
for i, d in enumerate([
    "Version fonctionnelle conçue, assemblée et testée",
    "Deux PCB routés, soudés et validés",
    "Firmware complet : acquisition, SD, Wi-Fi",
    "Design mécanique modulaire validé en conditions réelles",
    "Algorithme de vision  (F1 = 0.935) opérationnel",
]):
    tx(s16, Cm(2.0), Cm(3.85) + i * Cm(0.95), Cm(13.5), Cm(0.85),
       "✓  " + d, size=14, color=WHITE)

# À faire
rect(s16, Cm(17.5), Cm(2.65), Cm(15.5), Cm(1.0), AMBER)
tx(s16, Cm(17.8), Cm(2.7), Cm(15.0), Cm(0.88),
   "Prochaines étapes", size=18, bold=True, color=NAVY)
for i, t in enumerate([
    "Intégration NMEA 2000 et CTD lors d'une sortie Sailowtech",
    "Quantifier l'effet de serre résiduel sous le plexiglas",
    "Améliorer l'algorithme de vision (ciel uniforme, contre-jour)",
]):
    tx(s16, Cm(17.8), Cm(3.85) + i * Cm(1.2), Cm(15.0), Cm(1.1),
       "→  " + t, size=14, color=RGBColor(200, 220, 240))

# Message final
rect(s16, Cm(1.5), H - Cm(2.6), Cm(31.3), Cm(2.1), RGBColor(14, 33, 58))
tx(s16, Cm(2.0), H - Cm(2.5), Cm(30.5), Cm(1.8),
   "Station compacte, modulaire, low-tech et open-source —\n"
   "prête à rejoindre les expéditions Sailowtech.",
   size=17, bold=True, color=TEAL)


# ════════════════════════════════════════════════════════════
# SAUVEGARDE
# ════════════════════════════════════════════════════════════
out = r"C:\Users\romai\Desktop\saillowtech\Report_Sailowtech\presentation_sailowtech.pptx"
prs.save(out)
print(f"Fichier sauvegardé : {out}")
print(f"Nombre de slides   : {len(prs.slides)}")
